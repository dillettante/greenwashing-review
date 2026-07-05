from __future__ import annotations

import hashlib
import html
import re
import shutil
import subprocess
import tempfile
from pathlib import Path

from .models import SourcePage


SUPPORTED_EXTENSIONS = {".txt", ".md", ".html", ".htm", ".pdf", ".docx", ".pptx", ".png", ".jpg", ".jpeg"}


def _sanitize_text(text: str) -> str:
    """DOCX/XLSX(XML)에서 허용되지 않는 PDF 제어문자를 제거한다."""
    return "".join(
        char for char in text
        if char in "\t\n\r" or ("\x20" <= char <= "\ud7ff") or ("\ue000" <= char <= "\ufffd")
    )


def file_sha256(path: Path) -> str:
    digest = hashlib.sha256()
    with path.open("rb") as handle:
        for chunk in iter(lambda: handle.read(1024 * 1024), b""):
            digest.update(chunk)
    return digest.hexdigest()


def _plain_text(path: Path) -> list[str]:
    text = path.read_text(encoding="utf-8", errors="replace")
    if path.suffix.lower() in {".html", ".htm"}:
        text = re.sub(r"<(script|style).*?</\1>", " ", text, flags=re.I | re.S)
        text = re.sub(r"<[^>]+>", " ", text)
        text = html.unescape(text)
    return [text]


def _pdf(path: Path) -> list[str]:
    from pypdf import PdfReader

    return [(page.extract_text() or "") for page in PdfReader(str(path)).pages]


def _docx(path: Path) -> list[str]:
    from docx import Document

    doc = Document(str(path))
    chunks = [p.text for p in doc.paragraphs if p.text.strip()]
    for table in doc.tables:
        for row in table.rows:
            chunks.append(" | ".join(cell.text.strip() for cell in row.cells))
    return ["\n".join(chunks)]


def _pptx(path: Path) -> list[str]:
    from pptx import Presentation

    prs = Presentation(str(path))
    pages: list[str] = []
    for slide in prs.slides:
        chunks: list[str] = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                chunks.append(shape.text)
        pages.append("\n".join(chunks))
    return pages


def _image(path: Path) -> list[str]:
    if not shutil.which("tesseract"):
        return [""]
    with tempfile.TemporaryDirectory() as tmp:
        out = Path(tmp) / "ocr"
        subprocess.run(
            ["tesseract", str(path), str(out), "-l", "kor+eng"],
            check=False,
            capture_output=True,
            text=True,
            timeout=120,
        )
        text_path = out.with_suffix(".txt")
        return [text_path.read_text(encoding="utf-8", errors="replace") if text_path.exists() else ""]


def extract_file(path: Path, source_type: str) -> list[SourcePage]:
    suffix = path.suffix.lower()
    if suffix not in SUPPORTED_EXTENSIONS:
        return []
    extractors = {
        ".txt": _plain_text,
        ".md": _plain_text,
        ".html": _plain_text,
        ".htm": _plain_text,
        ".pdf": _pdf,
        ".docx": _docx,
        ".pptx": _pptx,
        ".png": _image,
        ".jpg": _image,
        ".jpeg": _image,
    }
    sha = file_sha256(path)
    document_id = sha[:12]
    return [
        SourcePage(
            document_id=document_id,
            filename=path.name,
            page=index,
            text=_sanitize_text(text).strip(),
            sha256=sha,
            source_type=source_type,
        )
        for index, text in enumerate(extractors[suffix](path), 1)
    ]


def extract_directory(directory: Path, source_type: str) -> tuple[list[SourcePage], list[str]]:
    pages: list[SourcePage] = []
    warnings: list[str] = []
    if not directory.exists():
        return pages, warnings
    for path in sorted(p for p in directory.rglob("*") if p.is_file()):
        if path.name.startswith("."):
            continue
        if path.suffix.lower() not in SUPPORTED_EXTENSIONS:
            warnings.append(f"지원하지 않는 파일 형식: {path.name}")
            continue
        try:
            extracted = extract_file(path, source_type)
        except Exception as exc:  # 파일별 실패를 격리한다.
            warnings.append(f"추출 실패 {path.name}: {exc}")
            continue
        if not extracted or not any(page.text.strip() for page in extracted):
            warnings.append(f"[확인 필요] 텍스트가 추출되지 않음: {path.name}")
        pages.extend(extracted)
    return pages, warnings
